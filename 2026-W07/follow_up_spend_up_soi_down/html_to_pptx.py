#!/usr/bin/env python3
"""Generate polished PPTX from source CSV data using Moloco template.

Pulls data directly from the same CSV as generate_report.py,
re-renders charts at slide-optimized DPI, and builds well-formatted tables.
"""

import os
import re
import base64
from io import BytesIO
from copy import deepcopy

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── Paths ────────────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_PATH = "/Users/yiran.shi/Desktop/2026_playground_moloco_center/claude-bq-agent/tmp/data/20260317_170128_1a36.csv"
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "../../Moloco Sales Deck template 2025.pptx")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "p0_investigate_deep_dive.pptx")

# ── Constants ────────────────────────────────────────────────────────────────
PERIOD_ORDER = ["1_Benchmark", "2_Jan", "3_Feb", "4_Mar_MTD"]
PERIOD_LABELS = {"1_Benchmark": "Benchmark\n(Dec 25-31)", "2_Jan": "Jan", "3_Feb": "Feb", "4_Mar_MTD": "Mar MTD"}
PERIOD_SHORT = {"1_Benchmark": "BM", "2_Jan": "Jan", "3_Feb": "Feb", "4_Mar_MTD": "Mar"}
CJK_APP_NAMES = {"TikTok ティックトック": "TikTok (JP)"}

# ── Colors ───────────────────────────────────────────────────────────────────
C_BLUE = "#4C9AFF"
C_RED = "#FF5252"
C_DARK = "#23272F"
C_GRAY = "#6B7280"
C_LIGHT_BG = "#F8FAFC"

RGB_BLUE = RGBColor(0x4C, 0x9A, 0xFF)
RGB_DARK = RGBColor(0x23, 0x27, 0x2F)
RGB_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RGB_HEADER_BG = RGBColor(0xE8, 0xF0, 0xFE)
RGB_ROW_ALT = RGBColor(0xF8, 0xFA, 0xFC)
RGB_HIGHLIGHT = RGBColor(0xFF, 0xFD, 0xE7)
RGB_GREEN = RGBColor(0x16, 0x83, 0x3A)
RGB_RED_TEXT = RGBColor(0xDC, 0x26, 0x26)
RGB_DIM = RGBColor(0x9C, 0xA3, 0xAF)
RGB_NOTE_BG = RGBColor(0xF1, 0xF5, 0xF9)
RGB_BORDER = RGBColor(0xE2, 0xE8, 0xF0)


# ── Helpers ──────────────────────────────────────────────────────────────────

def sanitize(name):
    return CJK_APP_NAMES.get(name, name)

def fmt_spend(v):
    if abs(v) >= 1_000_000: return f"${v/1_000_000:.1f}M"
    if abs(v) >= 1_000: return f"${v/1_000:.1f}K"
    return f"${v:.0f}"

def fmt_spend_table(v):
    return f"${v:,.0f}"

def fmt_pct(v):
    return f"{v:.2f}%"

def fmt_int(v):
    return f"{int(v):,}"

def calc_soi(mol, non_mol):
    total = mol + non_mol
    return (mol / total * 100) if total > 0 else 0.0


# ── Chart generation (high-res for slides) ───────────────────────────────────

def chart_to_bytes(fig, dpi=200):
    """Convert matplotlib figure to PNG bytes."""
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


def make_account_chart(df_total, title):
    """Account-level dual-axis bar+line chart. Returns BytesIO."""
    fig, ax1 = plt.subplots(figsize=(9, 3.2))
    fig.patch.set_facecolor("white")

    periods = [p for p in PERIOD_ORDER if p in df_total["period"].values]
    sub = df_total.set_index("period").loc[periods]
    labels = [PERIOD_LABELS.get(p, p) for p in periods]
    x = np.arange(len(periods))

    bars = ax1.bar(x, sub["avg_daily_spend"], color=C_BLUE, alpha=0.75, width=0.5,
                   edgecolor="white", linewidth=0.5)
    ax1.set_ylabel("Avg Daily Spend", fontsize=10, color=C_BLUE, fontweight="600")
    ax1.tick_params(axis="y", labelcolor=C_BLUE, labelsize=9)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: fmt_spend(v)))
    for i, bar in enumerate(bars):
        ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height(),
                 fmt_spend(sub["avg_daily_spend"].iloc[i]),
                 ha="center", va="bottom", fontsize=9, color=C_BLUE, fontweight="600")

    ax2 = ax1.twinx()
    ax2.plot(x, sub["soi_pct"], color=C_RED, marker="o", linewidth=2.5,
             markersize=9, zorder=5, markerfacecolor="white", markeredgewidth=2.5,
             markeredgecolor=C_RED)
    ax2.set_ylabel("SOI (%)", fontsize=10, color=C_RED, fontweight="600")
    ax2.tick_params(axis="y", labelcolor=C_RED, labelsize=9)
    for i, v in enumerate(sub["soi_pct"]):
        ax2.annotate(f"{v:.2f}%", (i, v), textcoords="offset points", xytext=(0, 12),
                     ha="center", fontsize=10, fontweight="bold", color=C_RED)

    ax1.set_xticks(x)
    ax1.set_xticklabels(labels, fontsize=10, fontweight="500")
    ax1.set_ylim(bottom=0)
    soi_vals = sub["soi_pct"]
    margin = max((soi_vals.max() - soi_vals.min()) * 0.5, 0.5)
    ax2.set_ylim(max(0, soi_vals.min() - margin), soi_vals.max() + margin * 2.5)
    ax1.grid(axis="y", alpha=0.15, linewidth=0.5)
    ax1.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)
    ax1.set_title(title, fontsize=13, fontweight="bold", color=C_DARK, pad=12)
    plt.tight_layout()
    return chart_to_bytes(fig)


def make_app_grid_chart(df_apps, account_name):
    """App-level grid of dual-axis charts. Returns BytesIO or None."""
    apps_total = df_apps.groupby("app_name")["avg_daily_spend"].sum().sort_values(ascending=False)
    top_apps = apps_total.head(8).index.tolist()
    df_top = df_apps[df_apps["app_name"].isin(top_apps)].copy()
    if df_top.empty:
        return None

    n = len(top_apps)
    ncols = min(4, n)
    nrows = (n + ncols - 1) // ncols
    fig, axes = plt.subplots(nrows, ncols, figsize=(3.5 * ncols, 2.8 * nrows))
    fig.patch.set_facecolor("white")
    if nrows == 1 and ncols == 1: axes = np.array([[axes]])
    elif nrows == 1: axes = axes.reshape(1, -1)
    elif ncols == 1: axes = axes.reshape(-1, 1)

    for idx, app in enumerate(top_apps):
        row, col = divmod(idx, ncols)
        ax1 = axes[row][col]
        grp = df_top[df_top["app_name"] == app]
        periods = [p for p in PERIOD_ORDER if p in grp["period"].values]
        if not periods:
            ax1.set_visible(False)
            continue
        sub = grp.set_index("period").loc[periods]
        x = np.arange(len(periods))

        ax1.bar(x, sub["avg_daily_spend"], color=C_BLUE, alpha=0.65, width=0.45,
                edgecolor="white", linewidth=0.3)
        ax1.set_ylabel("DRR", fontsize=7, color=C_BLUE)
        ax1.tick_params(axis="y", labelcolor=C_BLUE, labelsize=6)
        ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: fmt_spend(v)))
        ax1.set_ylim(bottom=0)

        ax2 = ax1.twinx()
        ax2.plot(x, sub["soi_pct"], color=C_RED, marker="o", linewidth=2,
                 markersize=5, zorder=5, markerfacecolor="white",
                 markeredgewidth=2, markeredgecolor=C_RED)
        ax2.set_ylabel("SOI%", fontsize=7, color=C_RED)
        ax2.tick_params(axis="y", labelcolor=C_RED, labelsize=6)
        for i, v in enumerate(sub["soi_pct"]):
            ax2.annotate(f"{v:.1f}%", (i, v), textcoords="offset points",
                         xytext=(0, 8), ha="center", fontsize=7, fontweight="bold",
                         color=C_RED)
        soi_vals = sub["soi_pct"]
        margin = max((soi_vals.max() - soi_vals.min()) * 0.5, 0.5)
        ax2.set_ylim(max(0, soi_vals.min() - margin), soi_vals.max() + margin * 2)

        ax1.set_xticks(x)
        ax1.set_xticklabels([PERIOD_SHORT.get(p, p) for p in periods], fontsize=7)
        ax1.set_title(sanitize(app)[:30], fontsize=8, fontweight="bold", color=C_DARK)
        ax1.grid(axis="y", alpha=0.12, linewidth=0.3)
        ax1.spines["top"].set_visible(False)
        ax2.spines["top"].set_visible(False)

    for idx in range(n, nrows * ncols):
        r, c = divmod(idx, ncols)
        axes[r][c].set_visible(False)

    fig.suptitle(f"{account_name} — App SOI vs Spend", fontsize=12,
                 fontweight="bold", color=C_DARK, y=1.02)
    plt.tight_layout()
    return chart_to_bytes(fig)


# ── Data computations ────────────────────────────────────────────────────────

def compute_impact(df_apps, df_total):
    """Compute SOI impact table data. Returns list of dicts."""
    bm = df_apps[df_apps["period"] == "1_Benchmark"]
    mar = df_apps[df_apps["period"] == "4_Mar_MTD"]
    bm_t = df_total[df_total["period"] == "1_Benchmark"]
    mar_t = df_total[df_total["period"] == "4_Mar_MTD"]
    if bm_t.empty or mar_t.empty:
        return []

    bm_denom = bm_t.iloc[0]["moloco_installs"] + bm_t.iloc[0]["non_moloco_installs"]
    mar_denom = mar_t.iloc[0]["moloco_installs"] + mar_t.iloc[0]["non_moloco_installs"]
    if bm_denom == 0 or mar_denom == 0:
        return []

    mar_total_spend = mar["avg_daily_spend"].sum()
    all_apps = set(bm["app_name"].tolist() + mar["app_name"].tolist())
    impacts = []
    for app in all_apps:
        bm_mol = bm.loc[bm["app_name"] == app, "moloco_installs"].sum()
        mar_mol = mar.loc[mar["app_name"] == app, "moloco_installs"].sum()
        bm_nonmol = bm.loc[bm["app_name"] == app, "non_moloco_installs"].sum()
        mar_nonmol = mar.loc[mar["app_name"] == app, "non_moloco_installs"].sum()
        bm_soi = calc_soi(bm_mol, bm_nonmol)
        mar_soi = calc_soi(mar_mol, mar_nonmol)
        mar_spend = mar.loc[mar["app_name"] == app, "avg_daily_spend"].sum()
        spend_share = (mar_spend / mar_total_spend * 100) if mar_total_spend > 0 else 0
        bm_c = (bm_mol / bm_denom * 100) if bm_denom > 0 else 0
        mar_c = (mar_mol / mar_denom * 100) if mar_denom > 0 else 0
        impacts.append({
            "app": sanitize(app), "bm_soi": bm_soi, "mar_soi": mar_soi,
            "soi_chg": mar_soi - bm_soi, "mar_spend": mar_spend,
            "spend_share": spend_share, "impact": mar_c - bm_c,
            "is_top": spend_share >= 10,
        })
    impacts.sort(key=lambda x: x["impact"], reverse=True)
    return impacts


def compute_app_pivot(df_apps):
    """Compute app pivot data. Returns list of dicts with per-period values."""
    apps_by_spend = df_apps.groupby("app_name")["avg_daily_spend"].sum().sort_values(ascending=False)
    rows = []
    for app_name in apps_by_spend.index:
        app_data = df_apps[df_apps["app_name"] == app_name].set_index("period")
        row = {"app": sanitize(app_name)}
        for p in PERIOD_ORDER:
            key = PERIOD_SHORT[p]
            if p in app_data.index:
                r = app_data.loc[p]
                row[f"{key}_drr"] = fmt_spend_table(r["avg_daily_spend"])
                row[f"{key}_soi"] = fmt_pct(r["soi_pct"])
            else:
                row[f"{key}_drr"] = "—"
                row[f"{key}_soi"] = "—"
        rows.append(row)
    return rows


# ── Slide helpers ────────────────────────────────────────────────────────────

def delete_all_slides(prs):
    """Remove all template slides."""
    ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    while len(prs.slides) > 0:
        elem = prs.slides._sldIdLst[0]
        rId = elem.get(ns + "id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
        prs.slides._sldIdLst.remove(elem)


def add_slide(prs):
    """Add a blank slide (layout 22) with no placeholders."""
    slide = prs.slides.add_slide(prs.slide_layouts[22])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def _set_cell_border(cell, color_rgb, width_pt=0.5, sides=("bottom",)):
    """Set thin borders on a table cell."""
    tc = cell._tc
    tcPr = tc.tcPr
    if tcPr is None:
        tcPr = etree.SubElement(tc, "{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr")

    border_map = {
        "bottom": "a:lnB", "top": "a:lnT", "left": "a:lnL", "right": "a:lnR"
    }
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for side in sides:
        tag = border_map[side]
        ln = etree.SubElement(tcPr, f"{{{ns}}}{tag.split(':')[1]}")
        ln.set("w", str(int(width_pt * 12700)))
        ln.set("cmpd", "sng")
        solidFill = etree.SubElement(ln, f"{{{ns}}}solidFill")
        srgb = etree.SubElement(solidFill, f"{{{ns}}}srgbClr")
        srgb.set("val", str(color_rgb))


def _no_border(cell):
    """Remove all borders from cell."""
    tc = cell._tc
    tcPr = tc.tcPr
    if tcPr is None:
        tcPr = etree.SubElement(tc, "{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr")
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for tag in ["lnB", "lnT", "lnL", "lnR"]:
        ln = etree.SubElement(tcPr, f"{{{ns}}}{tag}")
        ln.set("w", "0")
        noFill = etree.SubElement(ln, f"{{{ns}}}noFill")


def textbox(slide, left, top, width, height, text, size=10, bold=False,
            color=RGB_DARK, align=PP_ALIGN.LEFT, name="Calibri"):
    """Add a simple text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return box


def section_header_bar(slide):
    """Draw thin blue accent line at top."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.22), Inches(9.2), Pt(2.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGB_BLUE
    line.line.fill.background()


def slide_header(slide, title, subtitle=None):
    """Add header area: accent line + title + optional subtitle."""
    section_header_bar(slide)
    textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.35),
            title, size=16, bold=True, color=RGB_DARK)
    y = Inches(0.68)
    if subtitle:
        textbox(slide, Inches(0.4), Inches(0.62), Inches(9.2), Inches(0.25),
                subtitle, size=10, bold=False, color=RGBColor(0x64, 0x74, 0x8B))
        y = Inches(0.9)
    return y


def add_table(slide, headers, rows, left, top, width, available_h,
              col_widths=None, font_size=8, row_height_pt=18):
    """Add a formatted table. rows = list of list of (text, color_override|None).
    headers = list of (text, alignment).
    """
    n_cols = len(headers)
    n_rows = len(rows) + 1
    rh = Pt(row_height_pt)
    table_h = min(rh * n_rows, available_h)

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    tbl = shape.table
    # Disable banding
    tbl_pr = tbl._tbl.tblPr
    tbl_pr.set("bandRow", "0")
    tbl_pr.set("bandCol", "0")
    tbl_pr.set("firstRow", "0")
    tbl_pr.set("lastRow", "0")

    # Column widths
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = w
    else:
        for j in range(n_cols):
            tbl.columns[j].width = width // n_cols

    # Header row
    for j, (hdr_text, hdr_align) in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = hdr_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGB_HEADER_BG
        _set_cell_border(cell, RGB_BORDER, 0.75, ("bottom",))
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = RGB_DARK
            p.font.name = "Calibri"
            p.alignment = hdr_align
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(4)
        cell.margin_right = Pt(4)
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)

    # Data rows
    for i, row_data in enumerate(rows):
        is_alt = (i % 2 == 1)
        highlight = row_data.get("_highlight", False) if isinstance(row_data, dict) else False
        cells = row_data.get("cells", row_data) if isinstance(row_data, dict) else row_data

        for j, (text, color_override) in enumerate(cells):
            cell = tbl.cell(i + 1, j)
            cell.text = str(text)

            if highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGB_HIGHLIGHT
            elif is_alt:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGB_ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGB_WHITE

            _set_cell_border(cell, RGB_BORDER, 0.3, ("bottom",))

            _, hdr_align = headers[j]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Calibri"
                p.font.bold = False
                p.alignment = hdr_align
                if color_override == "green":
                    p.font.color.rgb = RGB_GREEN
                elif color_override == "red":
                    p.font.color.rgb = RGB_RED_TEXT
                elif color_override == "dim":
                    p.font.color.rgb = RGB_DIM
                elif color_override == "bold":
                    p.font.bold = True
                    p.font.color.rgb = RGB_DARK
                else:
                    p.font.color.rgb = RGB_DARK

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)

    return shape


def note_box(slide, left, top, width, text, font_size=8):
    """Styled note callout with blue left accent."""
    height = Pt(font_size * 4.5)
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGB_NOTE_BG
    bg.line.color.rgb = RGB_BORDER
    bg.line.width = Pt(0.5)
    # Accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(3), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGB_BLUE
    accent.line.fill.background()
    # Text
    textbox(slide, left + Pt(10), top + Pt(2), width - Pt(16), height - Pt(4),
            text, size=font_size, color=RGBColor(0x47, 0x55, 0x69))
    return height


# ── Build slides ─────────────────────────────────────────────────────────────

def build(prs, df):
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    ML = Inches(0.4)       # margin left
    CW = Inches(9.2)       # content width

    totals = df[df["app_name"] == "--- TOTAL ---"].copy()
    mar_totals = totals[totals["period"] == "4_Mar_MTD"].sort_values(
        "avg_daily_spend", ascending=False)
    account_order = mar_totals["account"].tolist()
    for a in df["account"].unique():
        if a not in account_order:
            account_order.append(a)

    # ════════════════════════════════════════════
    # TITLE SLIDE
    # ════════════════════════════════════════════
    slide = add_slide(prs)
    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGB_BLUE
    bar.line.fill.background()
    # Bottom accent
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Pt(4), SLIDE_W, Pt(4))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = RGB_BLUE
    bar2.line.fill.background()

    textbox(slide, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.8),
            "P0 Investigate: Spend Up + SOI Down", size=28, bold=True,
            color=RGB_DARK, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.5),
            "Deep Dive Analysis", size=16, bold=False,
            color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.8), Inches(3.0), Inches(8.4), Inches(0.8),
            "Benchmark (Dec 25-31) vs Jan / Feb / Mar MTD\n"
            "Criteria: Avg daily spend UP but SOI DOWN in March vs Benchmark",
            size=10, bold=False, color=RGBColor(0x94, 0xA3, 0xB8),
            align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # PER-ACCOUNT SLIDES
    # ════════════════════════════════════════════
    for acct in account_order:
        acct_total = df[(df["account"] == acct) & (df["app_name"] == "--- TOTAL ---")]
        acct_apps = df[(df["account"] == acct) & (df["app_name"] != "--- TOTAL ---")]

        # ──── Section divider ────
        slide = add_slide(prs)
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.8), SLIDE_W, Inches(2.0))
        band.fill.solid()
        band.fill.fore_color.rgb = RGB_BLUE
        band.line.fill.background()
        textbox(slide, Inches(0.8), Inches(2.1), Inches(8.4), Inches(0.7),
                acct, size=30, bold=True, color=RGB_WHITE, align=PP_ALIGN.CENTER)
        textbox(slide, Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.4),
                "Spend Up + SOI Down Analysis", size=13, bold=False,
                color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

        # ──── Slide: Trend Overview (chart + table) ────
        if not acct_total.empty:
            slide = add_slide(prs)
            y = slide_header(slide, acct, "Spend & SOI Trend Overview")

            chart_buf = make_account_chart(acct_total, f"{acct} — Daily Spend vs SOI")
            pic = slide.shapes.add_picture(chart_buf, Inches(0.5), y, Inches(9.0), Inches(2.3))

            y_table = y + Inches(2.4)
            headers = [
                ("Period", PP_ALIGN.LEFT),
                ("Total Spend", PP_ALIGN.RIGHT),
                ("Avg Daily Spend", PP_ALIGN.RIGHT),
                ("Moloco Installs", PP_ALIGN.RIGHT),
                ("Non-Moloco Installs", PP_ALIGN.RIGHT),
                ("SOI", PP_ALIGN.RIGHT),
            ]
            rows = []
            for p in PERIOD_ORDER:
                sub = acct_total[acct_total["period"] == p]
                if sub.empty:
                    continue
                r = sub.iloc[0]
                rows.append([
                    (PERIOD_SHORT[p], None),
                    (fmt_spend_table(r["total_spend"]), None),
                    (fmt_spend_table(r["avg_daily_spend"]), None),
                    (fmt_int(r["moloco_installs"]), None),
                    (fmt_int(r["non_moloco_installs"]), None),
                    (fmt_pct(r["soi_pct"]), None),
                ])

            col_w = [int(CW * f) for f in [0.12, 0.18, 0.18, 0.18, 0.20, 0.14]]
            add_table(slide, headers, rows, ML, y_table, CW,
                      SLIDE_H - y_table - Inches(0.2), col_widths=col_w,
                      font_size=9, row_height_pt=20)

        # ──── Slide: SOI Impact by App ────
        if not acct_apps.empty and not acct_total.empty:
            impacts = compute_impact(acct_apps, acct_total)
            if impacts:
                slide = add_slide(prs)
                y = slide_header(slide, acct, "SOI Impact by App (Benchmark → Mar)")

                note_h = note_box(
                    slide, ML, y, CW,
                    "SOI Impact = change in app's Moloco install share of total installs. "
                    "Negative = unattributed installs growing faster. "
                    "Highlighted = top spend apps (>=10% share).",
                    font_size=7)
                y += note_h + Inches(0.1)

                headers = [
                    ("App", PP_ALIGN.LEFT),
                    ("BM SOI", PP_ALIGN.RIGHT),
                    ("Mar SOI", PP_ALIGN.RIGHT),
                    ("SOI Chg", PP_ALIGN.RIGHT),
                    ("Mar DRR", PP_ALIGN.RIGHT),
                    ("Spend %", PP_ALIGN.RIGHT),
                    ("Impact (pp)", PP_ALIGN.RIGHT),
                ]
                rows = []
                for imp in impacts:
                    chg_c = "green" if imp["soi_chg"] >= 0 else "red"
                    imp_c = "green" if imp["impact"] >= 0 else "red"
                    chg_s = f"+{imp['soi_chg']:.2f}%" if imp["soi_chg"] >= 0 else f"{imp['soi_chg']:.2f}%"
                    imp_s = f"+{imp['impact']:.3f}" if imp["impact"] >= 0 else f"{imp['impact']:.3f}"
                    row = {
                        "_highlight": imp["is_top"],
                        "cells": [
                            (imp["app"][:30], "bold" if imp["is_top"] else None),
                            ("—" if imp["bm_soi"] == 0 else fmt_pct(imp["bm_soi"]), "dim" if imp["bm_soi"] == 0 else None),
                            ("—" if imp["mar_soi"] == 0 else fmt_pct(imp["mar_soi"]), "dim" if imp["mar_soi"] == 0 else None),
                            (chg_s, chg_c),
                            ("—" if imp["mar_spend"] == 0 else fmt_spend_table(imp["mar_spend"]), "dim" if imp["mar_spend"] == 0 else None),
                            ("—" if imp["spend_share"] == 0 else f"{imp['spend_share']:.1f}%", None),
                            (imp_s, imp_c),
                        ]
                    }
                    rows.append(row)

                col_w = [int(CW * f) for f in [0.22, 0.11, 0.11, 0.12, 0.16, 0.12, 0.16]]
                font = 8 if len(rows) <= 8 else 7
                rh = 18 if len(rows) <= 8 else 15
                add_table(slide, headers, rows, ML, y, CW,
                          SLIDE_H - y - Inches(0.15), col_widths=col_w,
                          font_size=font, row_height_pt=rh)

        # ──── Slide: App-Level SOI vs Spend (grid chart) ────
        if not acct_apps.empty:
            chart_buf = make_app_grid_chart(acct_apps, acct)
            if chart_buf:
                slide = add_slide(prs)
                y = slide_header(slide, acct, "App-Level SOI vs Spend")
                ch = Inches(4.2)
                cw = Inches(9.0)
                cl = (SLIDE_W - cw) // 2
                slide.shapes.add_picture(chart_buf, cl, y, cw, ch)

        # ──── Slide: App-Level Data (pivot table) ────
        if not acct_apps.empty:
            pivot = compute_app_pivot(acct_apps)
            if pivot:
                slide = add_slide(prs)
                y = slide_header(slide, acct, "App-Level Data: Benchmark vs Monthly Trend")

                headers = [("App", PP_ALIGN.LEFT)]
                for p in PERIOD_ORDER:
                    lbl = PERIOD_SHORT[p]
                    headers.append((f"{lbl} DRR", PP_ALIGN.RIGHT))
                    headers.append((f"{lbl} SOI", PP_ALIGN.RIGHT))

                rows = []
                for r in pivot:
                    row = [(r["app"][:25], "bold")]
                    for p in PERIOD_ORDER:
                        k = PERIOD_SHORT[p]
                        drr = r[f"{k}_drr"]
                        soi = r[f"{k}_soi"]
                        row.append((drr, "dim" if drr == "—" else None))
                        row.append((soi, "dim" if soi == "—" else None))
                    rows.append(row)

                # Column widths: app name wider, then equal for data
                app_w = int(CW * 0.17)
                data_w = (int(CW) - app_w) // 8
                col_w = [app_w] + [data_w] * 8

                font = 7 if len(rows) <= 10 else 6
                rh = 16 if len(rows) <= 10 else 13
                add_table(slide, headers, rows, ML, y, CW,
                          SLIDE_H - y - Inches(0.15), col_widths=col_w,
                          font_size=font, row_height_pt=rh)


def main():
    df = pd.read_csv(DATA_PATH)
    prs = Presentation(TEMPLATE_PATH)
    delete_all_slides(prs)
    build(prs, df)
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
